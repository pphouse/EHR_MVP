#!/usr/bin/env python3
"""
PowerPointファイルからテキストを抽出するスクリプト
"""
import sys
import os

try:
    from pptx import Presentation
except ImportError:
    print("python-pptxがインストールされていません。")
    print("インストールするには: pip install python-pptx")
    sys.exit(1)

def extract_text_from_pptx(filename):
    """PowerPointファイルからテキストを抽出"""
    try:
        prs = Presentation(filename)
        
        print(f"スライド数: {len(prs.slides)}")
        print("="*50)
        
        for i, slide in enumerate(prs.slides, 1):
            print(f"\n### スライド {i} ###")
            
            # タイトルを抽出
            if slide.shapes.title:
                print(f"タイトル: {slide.shapes.title.text}")
            
            # 各シェイプからテキストを抽出
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    if shape != slide.shapes.title:  # タイトル以外
                        print(f"- {shape.text}")
            
            # ノートを抽出
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                print(f"\nノート: {slide.notes_slide.notes_text_frame.text}")
            
            print("-"*50)
            
    except Exception as e:
        print(f"エラーが発生しました: {e}")
        return False
    
    return True

if __name__ == "__main__":
    filename = "GENIAC-PRIZE電子カルテv8.pptx"
    
    if os.path.exists(filename):
        print(f"ファイル '{filename}' からテキストを抽出します...")
        extract_text_from_pptx(filename)
    else:
        print(f"ファイル '{filename}' が見つかりません。")